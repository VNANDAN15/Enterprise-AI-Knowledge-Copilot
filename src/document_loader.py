import os
import csv
from pypdf import PdfReader
import docx
from pptx import Presentation
import pandas as pd


class UniversalDocumentLoader:
    def __init__(self, file_path):
        self.file_path = file_path
        self.ext = os.path.splitext(file_path)[1].lower()

    def extract_pages(self):
        """
        Extract content structured by page, slide, or sheet.
        Returns a list of dicts: [{"page_number": int, "text": str}, ...]
        """
        if self.ext == '.pdf':
            return self._extract_pdf()
        elif self.ext in ['.docx', '.doc']:
            return self._extract_docx()
        elif self.ext in ['.pptx', '.ppt']:
            return self._extract_pptx()
        elif self.ext in ['.xlsx', '.xls']:
            return self._extract_excel()
        elif self.ext == '.csv':
            return self._extract_csv()
        elif self.ext in ['.txt', '.md']:
            return self._extract_text()
        else:
            raise ValueError(f"Unsupported file extension: {self.ext}")

    def _extract_pdf(self):
        reader = PdfReader(self.file_path)
        pages = []
        for idx, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            pages.append({
                "page_number": idx + 1,
                "text": page_text.strip()
            })
        return pages

    def _extract_docx(self):
        doc = docx.Document(self.file_path)
        full_text = []

        # Extract paragraphs
        for p in doc.paragraphs:
            if p.text.strip():
                full_text.append(p.text.strip())

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_txt = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_txt:
                    full_text.append(row_txt)

        combined_text = "\n".join(full_text)
        if not combined_text.strip():
            return []

        # Virtual chunking into 1000-character pages for docx
        page_size = 1000
        pages = []
        for i in range(0, len(combined_text), page_size):
            chunk = combined_text[i:i + page_size]
            pages.append({
                "page_number": (i // page_size) + 1,
                "text": chunk.strip()
            })
        return pages

    def _extract_pptx(self):
        prs = Presentation(self.file_path)
        pages = []

        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            txt = "\n".join(slide_text)
            pages.append({
                "page_number": idx + 1,
                "text": txt if txt else f"[Slide {idx + 1} Visual Content]"
            })
        return pages

    def _extract_excel(self):
        excel_file = pd.ExcelFile(self.file_path)
        pages = []

        for idx, sheet_name in enumerate(excel_file.sheet_names):
            df = excel_file.parse(sheet_name)
            sheet_txt = f"Sheet: {sheet_name}\n" + df.to_markdown(index=False)
            pages.append({
                "page_number": idx + 1,
                "text": sheet_txt
            })
        return pages

    def _extract_csv(self):
        df = pd.read_csv(self.file_path)
        csv_txt = f"CSV Data ({len(df)} rows):\n" + df.to_markdown(index=False)

        # Split into virtual 2000 character pages if CSV is large
        page_size = 2000
        pages = []
        for i in range(0, len(csv_txt), page_size):
            chunk = csv_txt[i:i + page_size]
            pages.append({
                "page_number": (i // page_size) + 1,
                "text": chunk.strip()
            })
        return pages

    def _extract_text(self):
        with open(self.file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        if not content.strip():
            return []

        page_size = 1200
        pages = []
        for i in range(0, len(content), page_size):
            chunk = content[i:i + page_size]
            pages.append({
                "page_number": (i // page_size) + 1,
                "text": chunk.strip()
            })
        return pages
